#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
四川能投智慧光电智慧产业生态合作方案 PPT生成器
将HTML页面转换为统一风格的PowerPoint演示文稿
"""

import os
import re
from pathlib import Path
from typing import List, Dict, Any, Optional
from dataclasses import dataclass
from datetime import datetime

# HTML解析相关
from bs4 import BeautifulSoup

# PPT生成相关
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame

# 日志记录
import logging

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('ppt_generation.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

@dataclass
class PageContent:
    """页面内容数据结构"""
    page_number: int
    title: str
    content: str
    page_type: str  # 'cover', 'toc', 'content', 'thank_you'
    elements: List[Dict[str, Any]]
    
class PPTStyleConfig:
    """PPT样式配置类"""
    
    # 配色方案 - 蓝绿渐变
    PRIMARY_BLUE = RGBColor(59, 130, 246)  # #3b82f6
    PRIMARY_GREEN = RGBColor(16, 185, 129)  # #10b981
    BACKGROUND_LIGHT = RGBColor(248, 250, 252)  # #f8fafc
    TEXT_DARK = RGBColor(30, 41, 59)  # #1e293b
    TEXT_LIGHT = RGBColor(100, 116, 139)  # #64748b
    WHITE = RGBColor(255, 255, 255)
    
    # 字体配置
    FONT_FAMILY = "Noto Sans SC"
    TITLE_FONT_SIZE = Pt(36)
    SUBTITLE_FONT_SIZE = Pt(24)
    CONTENT_FONT_SIZE = Pt(18)
    SMALL_FONT_SIZE = Pt(14)
    
    # 布局配置
    SLIDE_WIDTH = Inches(13.33)  # 16:9比例
    SLIDE_HEIGHT = Inches(7.5)
    MARGIN_LEFT = Inches(0.5)
    MARGIN_RIGHT = Inches(0.5)
    MARGIN_TOP = Inches(0.5)
    MARGIN_BOTTOM = Inches(0.5)

class HTMLParser:
    """HTML文件解析器"""
    
    def __init__(self):
        self.pages_content: List[PageContent] = []
        
    def parse_all_files(self, webpages_folder: Path) -> List[PageContent]:
        """解析所有HTML页面"""
        logger.info(f"开始解析HTML文件夹: {webpages_folder}")
        
        # 获取所有HTML文件并排序
        html_files = sorted(
            [f for f in webpages_folder.glob("page*.html")],
            key=lambda x: int(re.search(r'page(\d+)', x.name).group(1))
        )
        
        logger.info(f"找到 {len(html_files)} 个HTML文件")
        
        for html_file in html_files:
            try:
                page_content = self._parse_single_page(html_file)
                if page_content:
                    self.pages_content.append(page_content)
                    logger.info(f"成功解析: {html_file.name} - {page_content.title}")
            except Exception as e:
                logger.error(f"解析文件 {html_file.name} 时出错: {str(e)}")
                
        return self.pages_content
    
    def _parse_single_page(self, html_file: Path) -> Optional[PageContent]:
        """解析单个HTML页面"""
        with open(html_file, 'r', encoding='utf-8') as f:
            content = f.read()
            
        soup = BeautifulSoup(content, 'html.parser')
        
        # 提取页面编号
        page_number = int(re.search(r'page(\d+)', html_file.name).group(1))
        
        # 提取标题
        title = self._extract_title(soup)
        
        # 提取内容
        text_content = self._extract_content(soup)
        
        # 确定页面类型
        page_type = self._determine_page_type(page_number, title, text_content)
        
        # 提取页面元素
        elements = self._extract_elements(soup)
        
        return PageContent(
            page_number=page_number,
            title=title,
            content=text_content,
            page_type=page_type,
            elements=elements
        )
    
    def _extract_title(self, soup: BeautifulSoup) -> str:
        """提取页面标题"""
        # 首先尝试从HTML title标签获取
        title_tag = soup.find('title')
        if title_tag and title_tag.get_text(strip=True):
            return title_tag.get_text(strip=True)
        
        # 尝试多种标题提取方式
        title_selectors = [
            'h1', 'h2', '.text-4xl', '.text-5xl', '.text-6xl',
            '.title', '.main-title', '.page-title', '[class*="title"]'
        ]
        
        for selector in title_selectors:
            title_elem = soup.select_one(selector)
            if title_elem and title_elem.get_text(strip=True):
                title_text = title_elem.get_text(strip=True)
                # 对于封面页，可能有多个标题元素，需要合并
                if selector in ['h1', 'h2'] and len(title_text) < 50:
                    # 查找相邻的标题元素
                    next_title = title_elem.find_next(['h1', 'h2'])
                    if next_title and len(next_title.get_text(strip=True)) < 100:
                        return f"{title_text} {next_title.get_text(strip=True)}"
                return title_text
        
        # 如果没找到标题，尝试从页面内容中提取最显眼的文本
        large_text_elements = soup.find_all(attrs={'class': re.compile(r'text-(4xl|5xl|6xl|2xl|3xl)')}) 
        for elem in large_text_elements:
            text = elem.get_text(strip=True)
            if text and len(text) < 200:  # 标题通常较短
                return text
                
        return "无标题"
    
    def _extract_content(self, soup: BeautifulSoup) -> str:
        """提取页面内容"""
        # 移除script和style标签
        for script in soup(["script", "style"]):
            script.decompose()
        
        # 提取主要内容区域
        content_areas = []
        
        # 查找主要内容容器
        main_content = soup.find('div', class_=re.compile(r'(content|main|body)'))
        if not main_content:
            main_content = soup.find('body')
        
        if main_content:
            # 提取段落文本
            paragraphs = main_content.find_all(['p', 'div', 'span'], 
                                              class_=re.compile(r'text-(lg|xl|base|sm)'))
            for p in paragraphs:
                text = p.get_text(strip=True)
                if text and len(text) > 10:  # 过滤太短的文本
                    content_areas.append(text)
            
            # 如果没有找到段落，提取所有文本
            if not content_areas:
                text_content = main_content.get_text(separator='\n', strip=True)
                lines = [line.strip() for line in text_content.split('\n') 
                        if line.strip() and len(line.strip()) > 5]
                content_areas = lines[:20]  # 限制行数
        
        return '\n'.join(content_areas) if content_areas else "无内容"
    
    def _determine_page_type(self, page_number: int, title: str, content: str) -> str:
        """确定页面类型"""
        if page_number == 1:
            return 'cover'
        elif page_number == 2 or '目录' in title or 'contents' in title.lower():
            return 'toc'
        elif page_number == 30 or '谢谢' in title or 'thank' in title.lower():
            return 'thank_you'
        else:
            return 'content'
    
    def _extract_elements(self, soup: BeautifulSoup) -> List[Dict[str, Any]]:
        """提取页面元素（图片、列表等）"""
        elements = []
        
        # 提取列表项（包括目录项和功能点）
        list_items = []
        
        # 查找ul/ol列表
        for ul in soup.find_all(['ul', 'ol']):
            items = [li.get_text(strip=True) for li in ul.find_all('li') if li.get_text(strip=True)]
            if items:
                list_items.extend(items)
        
        # 查找带有特定class的列表项（如toc-item）
        toc_items = soup.find_all(class_=re.compile(r'(toc-item|list-item|feature)'))
        for item in toc_items:
            text = item.get_text(strip=True)
            if text and text not in list_items:
                list_items.append(text)
        
        if list_items:
            elements.append({
                'type': 'list',
                'items': list_items[:15]  # 限制数量
            })
        
        # 提取卡片式内容
        cards = soup.find_all(class_=re.compile(r'(card|feature-card|section)'))
        card_contents = []
        for card in cards:
            title_elem = card.find(['h2', 'h3', 'h4'])
            content_elem = card.find('p')
            if title_elem and content_elem:
                card_contents.append({
                    'title': title_elem.get_text(strip=True),
                    'content': content_elem.get_text(strip=True)
                })
        
        if card_contents:
            elements.append({
                'type': 'cards',
                'items': card_contents[:8]  # 限制数量
            })
        
        # 提取表格
        for table in soup.find_all('table'):
            rows = []
            for tr in table.find_all('tr'):
                cells = [td.get_text(strip=True) for td in tr.find_all(['td', 'th'])]
                if cells:
                    rows.append(cells)
            if rows:
                elements.append({
                    'type': 'table',
                    'rows': rows
                })
        
        return elements

class PPTGenerator:
    """PPT生成器"""
    
    def __init__(self):
        self.prs = Presentation()
        self.style_config = PPTStyleConfig()
        
        # 设置幻灯片尺寸为16:9
        self.prs.slide_width = self.style_config.SLIDE_WIDTH
        self.prs.slide_height = self.style_config.SLIDE_HEIGHT
        
    def generate_presentation(self, pages_content: List[PageContent], output_path: str):
        """生成完整的PPT演示文稿"""
        logger.info("开始生成PPT演示文稿")
        
        for page_content in pages_content:
            try:
                if page_content.page_type == 'cover':
                    self._create_cover_slide(page_content)
                elif page_content.page_type == 'toc':
                    self._create_toc_slide(page_content)
                elif page_content.page_type == 'thank_you':
                    self._create_thank_you_slide(page_content)
                else:
                    self._create_content_slide(page_content)
                    
                logger.info(f"已生成第 {page_content.page_number} 页: {page_content.title}")
                
            except Exception as e:
                logger.error(f"生成第 {page_content.page_number} 页时出错: {str(e)}")
        
        # 保存PPT文件
        self.prs.save(output_path)
        logger.info(f"PPT文件已保存至: {output_path}")
    
    def _add_gradient_background(self, slide: Slide):
        """为幻灯片添加渐变背景"""
        try:
            # 添加背景形状
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, 
                self.style_config.SLIDE_WIDTH, 
                self.style_config.SLIDE_HEIGHT
            )
            
            # 设置背景颜色为浅灰色
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self.style_config.BACKGROUND_LIGHT
            
            # 移到最底层
            background.element.getparent().remove(background.element)
            slide.shapes._spTree.insert(2, background.element)
            
        except Exception as e:
            logger.warning(f"添加背景时出错: {str(e)}")
    
    def _create_cover_slide(self, page_content: PageContent):
        """创建封面页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加背景
        self._add_gradient_background(slide)
        
        # 主标题
        title_box = slide.shapes.add_textbox(
            self.style_config.MARGIN_LEFT, 
            Inches(2), 
            Inches(10), 
            Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        
        # 分割标题
        title_parts = page_content.title.split('智慧产业生态合作方案')
        if len(title_parts) == 2:
            # 公司名称
            p1 = title_frame.paragraphs[0]
            p1.text = title_parts[0].strip()
            p1.font.name = self.style_config.FONT_FAMILY
            p1.font.size = Pt(32)
            p1.font.color.rgb = self.style_config.TEXT_DARK
            p1.alignment = PP_ALIGN.CENTER
            
            # 方案标题
            p2 = title_frame.add_paragraph()
            p2.text = "智慧产业生态合作方案"
            p2.font.name = self.style_config.FONT_FAMILY
            p2.font.size = Pt(40)
            p2.font.bold = True
            p2.font.color.rgb = self.style_config.PRIMARY_BLUE
            p2.alignment = PP_ALIGN.CENTER
        else:
            p1 = title_frame.paragraphs[0]
            p1.text = page_content.title
            p1.font.name = self.style_config.FONT_FAMILY
            p1.font.size = self.style_config.TITLE_FONT_SIZE
            p1.font.bold = True
            p1.font.color.rgb = self.style_config.PRIMARY_BLUE
            p1.alignment = PP_ALIGN.CENTER
        
        # 副标题/描述
        if page_content.elements:
            for element in page_content.elements:
                if element['type'] == 'list' and element['items']:
                    desc_box = slide.shapes.add_textbox(
                        Inches(2), Inches(4.5), Inches(8), Inches(2)
                    )
                    desc_frame = desc_box.text_frame
                    desc_frame.clear()
                    
                    for i, item in enumerate(element['items'][:3]):
                        if i == 0:
                            p = desc_frame.paragraphs[0]
                        else:
                            p = desc_frame.add_paragraph()
                        p.text = f"• {item}"
                        p.font.name = self.style_config.FONT_FAMILY
                        p.font.size = Pt(16)
                        p.font.color.rgb = self.style_config.TEXT_LIGHT
                        p.alignment = PP_ALIGN.CENTER
                    break
        
        # 日期
        date_box = slide.shapes.add_textbox(
            Inches(8), Inches(6), Inches(4), Inches(1)
        )
        date_frame = date_box.text_frame
        date_frame.clear()
        p = date_frame.paragraphs[0]
        p.text = "2025年08月20日"
        p.font.name = self.style_config.FONT_FAMILY
        p.font.size = Pt(14)
        p.font.color.rgb = self.style_config.TEXT_LIGHT
        p.alignment = PP_ALIGN.RIGHT
    
    def _create_toc_slide(self, page_content: PageContent):
        """创建目录页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加背景
        self._add_gradient_background(slide)
        
        # 标题
        title_box = slide.shapes.add_textbox(
            self.style_config.MARGIN_LEFT, 
            self.style_config.MARGIN_TOP, 
            Inches(8), 
            Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = "目录"
        p.font.name = self.style_config.FONT_FAMILY
        p.font.size = self.style_config.TITLE_FONT_SIZE
        p.font.bold = True
        p.font.color.rgb = self.style_config.TEXT_DARK
        
        # 目录内容
        if page_content.elements:
            for element in page_content.elements:
                if element['type'] == 'list':
                    content_box = slide.shapes.add_textbox(
                        Inches(1), Inches(1.5), Inches(10), Inches(5)
                    )
                    content_frame = content_box.text_frame
                    content_frame.clear()
                    
                    # 按章节组织目录项
                    sections = []
                    current_section = None
                    
                    for item in element['items']:
                        if any(x in item for x in ['一、', '二、', '三、', '四、', '五、', '六、']):
                            if current_section:
                                sections.append(current_section)
                            current_section = {'title': item, 'items': []}
                        elif current_section:
                            current_section['items'].append(item)
                    
                    if current_section:
                        sections.append(current_section)
                    
                    # 渲染目录
                    for i, section in enumerate(sections[:6]):
                        if i == 0:
                            p = content_frame.paragraphs[0]
                        else:
                            p = content_frame.add_paragraph()
                        
                        p.text = section['title']
                        p.font.name = self.style_config.FONT_FAMILY
                        p.font.size = Pt(18)
                        p.font.bold = True
                        p.font.color.rgb = self.style_config.PRIMARY_BLUE
                        p.space_after = Pt(6)
                        
                        # 添加子项
                        for sub_item in section['items'][:3]:
                            p_sub = content_frame.add_paragraph()
                            p_sub.text = f"  • {sub_item}"
                            p_sub.font.name = self.style_config.FONT_FAMILY
                            p_sub.font.size = Pt(14)
                            p_sub.font.color.rgb = self.style_config.TEXT_LIGHT
                            p_sub.space_after = Pt(3)
                    break
    
    def _create_content_slide(self, page_content: PageContent):
        """创建内容页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加背景
        self._add_gradient_background(slide)
        
        # 标题
        title_box = slide.shapes.add_textbox(
            self.style_config.MARGIN_LEFT, 
            self.style_config.MARGIN_TOP, 
            Inches(10), 
            Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = page_content.title
        p.font.name = self.style_config.FONT_FAMILY
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.style_config.TEXT_DARK
        
        # 内容区域
        content_top = Inches(1.5)
        
        # 处理不同类型的元素
        if page_content.elements:
            y_offset = content_top
            
            for element in page_content.elements:
                if element['type'] == 'list' and y_offset < Inches(6):
                    list_box = slide.shapes.add_textbox(
                        Inches(1), y_offset, Inches(10), Inches(4)
                    )
                    list_frame = list_box.text_frame
                    list_frame.clear()
                    
                    for i, item in enumerate(element['items'][:8]):
                        if i == 0:
                            p = list_frame.paragraphs[0]
                        else:
                            p = list_frame.add_paragraph()
                        p.text = f"• {item}"
                        p.font.name = self.style_config.FONT_FAMILY
                        p.font.size = Pt(16)
                        p.font.color.rgb = self.style_config.TEXT_DARK
                        p.space_after = Pt(6)
                    
                    y_offset += Inches(3)
                    
                elif element['type'] == 'cards' and y_offset < Inches(5):
                    # 卡片式布局
                    cards_per_row = min(3, len(element['items']))
                    card_width = Inches(10) / cards_per_row - Inches(0.2)
                    
                    for i, card in enumerate(element['items'][:6]):
                        if y_offset >= Inches(6):
                            break
                            
                        x_pos = Inches(1) + (i % cards_per_row) * (card_width + Inches(0.2))
                        if i > 0 and i % cards_per_row == 0:
                            y_offset += Inches(1.8)
                        
                        card_box = slide.shapes.add_textbox(
                            x_pos, y_offset, card_width, Inches(1.5)
                        )
                        card_frame = card_box.text_frame
                        card_frame.clear()
                        
                        # 卡片标题
                        p_title = card_frame.paragraphs[0]
                        p_title.text = card['title']
                        p_title.font.name = self.style_config.FONT_FAMILY
                        p_title.font.size = Pt(14)
                        p_title.font.bold = True
                        p_title.font.color.rgb = self.style_config.PRIMARY_BLUE
                        
                        # 卡片内容
                        p_content = card_frame.add_paragraph()
                        p_content.text = card['content'][:100] + "..." if len(card['content']) > 100 else card['content']
                        p_content.font.name = self.style_config.FONT_FAMILY
                        p_content.font.size = Pt(12)
                        p_content.font.color.rgb = self.style_config.TEXT_LIGHT
                    
                    y_offset += Inches(2)
        
        # 如果没有结构化元素，添加纯文本内容
        else:
            content_box = slide.shapes.add_textbox(
                Inches(1), content_top, Inches(10), Inches(4.5)
            )
            content_frame = content_box.text_frame
            content_frame.clear()
            
            # 分段显示内容
            content_lines = page_content.content.split('\n')[:15]
            for i, line in enumerate(content_lines):
                if line.strip():
                    if i == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                    p.text = line.strip()
                    p.font.name = self.style_config.FONT_FAMILY
                    p.font.size = Pt(14)
                    p.font.color.rgb = self.style_config.TEXT_DARK
                    p.space_after = Pt(6)
        
        # 页脚
        footer_box = slide.shapes.add_textbox(
            Inches(10), Inches(6.5), Inches(2), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_frame.clear()
        p = footer_frame.paragraphs[0]
        p.text = f"{page_content.page_number}/30"
        p.font.name = self.style_config.FONT_FAMILY
        p.font.size = Pt(12)
        p.font.color.rgb = self.style_config.TEXT_LIGHT
        p.alignment = PP_ALIGN.RIGHT
    
    def _create_thank_you_slide(self, page_content: PageContent):
        """创建感谢页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加背景
        self._add_gradient_background(slide)
        
        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(4), Inches(2.5), Inches(4), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = "谢谢"
        p.font.name = self.style_config.FONT_FAMILY
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.style_config.PRIMARY_BLUE
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(3), Inches(4), Inches(6), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        p = subtitle_frame.paragraphs[0]
        p.text = "感谢您的关注与支持"
        p.font.name = self.style_config.FONT_FAMILY
        p.font.size = Pt(20)
        p.font.color.rgb = self.style_config.TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        
        # 联系信息
        contact_box = slide.shapes.add_textbox(
            Inches(2), Inches(5), Inches(8), Inches(1.5)
        )
        contact_frame = contact_box.text_frame
        contact_frame.clear()
        
        contact_info = [
            "四川能投智慧光电有限公司",
            "电子科技大学",
            "合作共赢 · 智创未来"
        ]
        
        for i, info in enumerate(contact_info):
            if i == 0:
                p = contact_frame.paragraphs[0]
            else:
                p = contact_frame.add_paragraph()
            p.text = info
            p.font.name = self.style_config.FONT_FAMILY
            p.font.size = Pt(16)
            p.font.color.rgb = self.style_config.TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(8)

def main():
    """主函数"""
    # 配置路径
    webpages_folder = r"C:\Users\ARIS\Downloads\webpages"
    output_file = r"C:\Users\ARIS\Downloads\四川能投智慧光电智慧产业生态合作方案.pptx"
    
    try:
        # 检查输入文件夹是否存在
        if not os.path.exists(webpages_folder):
            logger.error(f"输入文件夹不存在: {webpages_folder}")
            return
        
        logger.info("=== 四川能投智慧光电智慧产业生态合作方案 PPT生成器 ===")
        logger.info(f"开始时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        
        # 解析HTML文件
        parser = HTMLParser()
        pages_content = parser.parse_all_files(Path(webpages_folder))
        
        if not pages_content:
            logger.error("未找到有效的HTML页面内容")
            return
        
        # 生成PPT
        generator = PPTGenerator()
        generator.generate_presentation(pages_content, output_file)
        
        logger.info(f"PPT生成完成！共处理 {len(pages_content)} 个页面")
        logger.info(f"输出文件: {output_file}")
        print(f"\n✅ 成功生成PPT文档: {output_file}")
        print(f"📄 共处理 {len(pages_content)} 个页面")
        print(f"🎨 使用统一的蓝绿渐变设计风格")
        print(f"📝 包含封面、目录、内容页和感谢页")
        
    except Exception as e:
        logger.error(f"程序执行出错: {str(e)}")
        print(f"\n❌ 生成PPT时发生错误: {str(e)}")
        raise

if __name__ == "__main__":
    main()